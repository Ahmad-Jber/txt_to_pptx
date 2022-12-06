
from pptx import Presentation
from textblob import TextBlob
from Summarizer import nltk_summarizer

self = Presentation()


def create_powerpoint(data, name):
    ppt_blob = TextBlob(data)
    for sentence in ppt_blob.sentences:
        slide = self.slides.add_slide(self.slide_layouts[1])
        for i in range(5):
            sent = str(sentence)
            slide.shapes.title.text = "Patent Cooperation Treaty"
            slide.placeholders[1].text = nltk_summarizer(sent)
    self.save(name)


def delete_slide(prs, slide):
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]